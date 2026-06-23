import io

from pptx import Presentation


def parse_slides(data: bytes) -> list[list[str]]:
    prs = Presentation(io.BytesIO(data))
    result = []

    for slide in prs.slides:
        texts = []
    
        for shape in slide.shapes:
    
            if shape.has_text_frame:
    
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
    
                    if text:
                        texts.append(text)
    
        result.append(texts)

    return result

'''
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pptx_parser.py <file.pptx>")
        sys.exit(1)

    with open(sys.argv[1], "rb") as f:
        slides = parse_slides(f.read())

    for i, slide in enumerate(slides):
        print(f"--- Slide {i + 1} ---")
        for text in slide:
            print(text)
'''
