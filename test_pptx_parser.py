import io
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_parser import parse_slides

PPTX_DIR = Path(__file__).parent / "pptx"
SAMPLE_FILES = list(PPTX_DIR.glob("*.pptx"))


def make_pptx(slides: list[list[str]]) -> bytes:
    """슬라이드별 텍스트 목록을 받아 pptx bytes를 생성."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for texts in slides:
        slide = prs.slides.add_slide(blank_layout)
        for i, text in enumerate(texts):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(8), Inches(0.5))
            txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── 단위 테스트 ──────────────────────────────────────────────────────────────

def test_single_slide_single_text():
    data = make_pptx([["안녕하세요"]])
    result = parse_slides(data)
    assert result == [["안녕하세요"]]


def test_single_slide_multiple_texts():
    data = make_pptx([["제목", "내용", "설명"]])
    result = parse_slides(data)
    assert result == [["제목", "내용", "설명"]]


def test_multiple_slides():
    data = make_pptx([["슬라이드 1"], ["슬라이드 2"], ["슬라이드 3"]])
    result = parse_slides(data)
    assert len(result) == 3
    assert result[0] == ["슬라이드 1"]
    assert result[1] == ["슬라이드 2"]
    assert result[2] == ["슬라이드 3"]


def test_empty_paragraphs_are_filtered():
    """공백 문자열만 있는 단락은 제거되어야 한다."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.text = "첫 번째"
    tf.add_paragraph().text = "   "   # 공백만 있는 단락
    tf.add_paragraph().text = "두 번째"
    buf = io.BytesIO()
    prs.save(buf)

    result = parse_slides(buf.getvalue())
    assert result == [["첫 번째", "두 번째"]]


def test_empty_pptx_returns_empty_list():
    prs = Presentation()
    buf = io.BytesIO()
    prs.save(buf)
    result = parse_slides(buf.getvalue())
    assert result == []


def test_returns_list_of_lists():
    data = make_pptx([["텍스트"]])
    result = parse_slides(data)
    assert isinstance(result, list)
    assert isinstance(result[0], list)


# ── 실제 파일 통합 테스트 ─────────────────────────────────────────────────────

@pytest.mark.parametrize("pptx_path", SAMPLE_FILES, ids=lambda p: p.name)
def test_real_file_parses_without_error(pptx_path: Path):
    """pptx/ 폴더의 실제 파일들이 오류 없이 파싱되어야 한다."""
    data = pptx_path.read_bytes()
    result = parse_slides(data)
    assert isinstance(result, list)
    assert all(isinstance(slide, list) for slide in result)


@pytest.mark.parametrize("pptx_path", SAMPLE_FILES, ids=lambda p: p.name)
def test_real_file_has_content(pptx_path: Path):
    """실제 파일은 최소 한 슬라이드에 텍스트가 있어야 한다."""
    data = pptx_path.read_bytes()
    result = parse_slides(data)
    non_empty_slides = [s for s in result if s]
    assert len(non_empty_slides) > 0, f"{pptx_path.name}: 텍스트가 있는 슬라이드가 없음"


@pytest.mark.parametrize("pptx_path", SAMPLE_FILES, ids=lambda p: p.name)
def test_real_file_texts_are_strings(pptx_path: Path):
    """파싱된 텍스트는 모두 str 타입이어야 한다."""
    data = pptx_path.read_bytes()
    result = parse_slides(data)
    for slide in result:
        for text in slide:
            assert isinstance(text, str)
            assert text.strip() == text, f"앞뒤 공백이 남아있음: {repr(text)}"
